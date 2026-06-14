"""Education course pack — remote education Skills + host-written Office artifacts.

Run:
    python examples/skills_executor/06_executable_education_course_pack.py

    The script also works when invoked by absolute path from outside the repo.

Environment:
    DEEPSEEK_API_KEY in the shell or .env file.
    Set DYNAMIC_TASK_MODEL_PROVIDER=ollama for local Ollama instead.
    Optional: pip install python-docx reportlab python-pptx openpyxl
    (each artifact is skipped gracefully if its library is missing).

Expected key output from a real DeepSeek run:
    skill status: success
    course title: <generated>
    artifacts written: up to 6 (docx/pdf/pptx/xlsx/json present-library dependent)

New-standard Skills model
-------------------------
This example uses real third-party Skills instead of a local demo Skill:
GarethManning's education Skill library for pedagogy plus Anthropic's artifact
Skills for document/spreadsheet/slide/PDF output guidance. The HOST still owns
every side effect — writing the .docx/.pdf/.pptx/.xlsx artifacts. Python package
dependencies are the host's responsibility, never a Skill's: a missing library
simply skips that artifact (we do not install packages at runtime).
"""

from __future__ import annotations

import asyncio
import json
import sys
import tempfile
from pathlib import Path
from typing import Any, cast

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from agently import Agently
from examples.dynamic_task._shared import configure_model

COURSE_BRIEF = "Design a B1 (intermediate) Business English course pack focused on running effective meetings."

REMOTE_SKILLS = [
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/curriculum-assessment/backwards-design-unit-planner", "trust_level": "remote"},
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/eal-language-development/language-demand-analyser", "trust_level": "remote"},
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/eal-language-development/vocabulary-tiering-tool", "trust_level": "remote"},
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/memory-learning-science/retrieval-practice-generator", "trust_level": "remote"},
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/memory-learning-science/spaced-practice-scheduler", "trust_level": "remote"},
    {"source": "GarethManning/education-agent-skills", "subpath": "skills/curriculum-assessment/formative-assessment-technique-selector", "trust_level": "remote"},
    {"source": "anthropics/skills", "subpath": "skills/docx", "trust_level": "remote"},
    {"source": "anthropics/skills", "subpath": "skills/pdf", "trust_level": "remote"},
    {"source": "anthropics/skills", "subpath": "skills/pptx", "trust_level": "remote"},
    {"source": "anthropics/skills", "subpath": "skills/xlsx", "trust_level": "remote"},
]


# ═══════════════════════════════════════════════════════════════════════════════
# HOST artifact writers — controlled, library-guarded (no runtime pip install)
# ═══════════════════════════════════════════════════════════════════════════════

def _write_docx(path: Path, title: str, sections: list[tuple[str, list[str]]]) -> bool:
    try:
        from docx import Document
    except ImportError:
        return False
    document = Document()
    document.add_heading(title, level=0)
    for heading, paragraphs in sections:
        document.add_heading(heading, level=1)
        for paragraph in paragraphs:
            document.add_paragraph(str(paragraph))
    document.save(str(path))
    return True


def _write_pdf(path: Path, title: str, paragraphs: list[str]) -> bool:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
        from xml.sax.saxutils import escape
    except ImportError:
        return False
    styles = getSampleStyleSheet()
    story: list[Any] = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for paragraph in paragraphs:
        story.append(Paragraph(escape(str(paragraph)), styles["BodyText"]))
        story.append(Spacer(1, 8))
    SimpleDocTemplate(str(path), pagesize=letter).build(story)
    return True


def _write_pptx(path: Path, slides: list[dict[str, Any]]) -> bool:
    try:
        from pptx import Presentation
    except ImportError:
        return False
    presentation = Presentation()
    for item in slides[:16]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        cast(Any, slide.shapes.title).text = str(item.get("title") or "Lesson")
        body = cast(Any, slide.placeholders[1]).text_frame
        body.clear()
        bullets = [str(b) for b in (item.get("bullets") or ["—"])][:6]
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
    presentation.save(str(path))
    return True


def _write_xlsx(path: Path, sheet_name: str, headers: list[str], rows: list[list[Any]]) -> bool:
    try:
        from openpyxl import Workbook
    except ImportError:
        return False
    workbook = Workbook()
    sheet = cast(Any, workbook.active)
    sheet.title = sheet_name[:31]
    sheet.append(headers)
    for row in rows:
        sheet.append(row)
    workbook.save(str(path))
    return True


def write_artifacts(out_dir: Path, course: dict[str, Any]) -> dict[str, str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    written: dict[str, str] = {}

    # Always-writable: the structured course plan as JSON.
    plan_path = out_dir / "course_plan.json"
    plan_path.write_text(json.dumps(course, ensure_ascii=False, indent=2), encoding="utf-8")
    written["course_plan"] = str(plan_path)

    sections = [(str(s.get("heading", "Section")), [str(p) for p in (s.get("points") or [])])
                for s in (course.get("teacher_guide_sections") or [])]
    if sections and _write_docx(out_dir / "teacher_guide.docx", str(course.get("course_title", "Course")), sections):
        written["teacher_guide"] = str(out_dir / "teacher_guide.docx")

    handout = [str(p) for p in (course.get("student_handout_paragraphs") or [])]
    if handout and _write_pdf(out_dir / "student_handout.pdf", "Student Handout", handout):
        written["student_handout"] = str(out_dir / "student_handout.pdf")

    slides = [s for s in (course.get("lesson_slides") or []) if isinstance(s, dict)]
    if slides and _write_pptx(out_dir / "lesson_slides.pptx", slides):
        written["lesson_slides"] = str(out_dir / "lesson_slides.pptx")

    vocab = [[str(v.get("term", "")), str(v.get("definition", ""))] for v in (course.get("vocabulary") or [])]
    if vocab and _write_xlsx(out_dir / "vocabulary_bank.xlsx", "Vocabulary", ["Term", "Definition"], vocab):
        written["vocabulary_bank"] = str(out_dir / "vocabulary_bank.xlsx")

    rubric = [(str(c.get("criterion", "Criterion")), [str(c.get("descriptor", ""))])
              for c in (course.get("rubric_criteria") or [])]
    if rubric and _write_docx(out_dir / "assessment_rubric.docx", "Assessment Rubric", rubric):
        written["assessment_rubric"] = str(out_dir / "assessment_rubric.docx")

    return written


async def main() -> None:
    provider = configure_model(temperature=0.3)
    print(f"Model provider: {provider}\n")

    Agently.skills_executor.configure(
        registry_root=str(ROOT / ".example_runtime" / "skills_executor" / "education_course_pack" / "registry"),
        allowed_trust_levels=["local", "remote"],
    )
    agent = Agently.create_agent("course-designer")
    agent.use_skills(REMOTE_SKILLS, mode="required")

    divider = "=" * 60
    print(divider)
    print("Education Course Pack — prompt-only Skill + host artifacts")
    print(f"Brief: {COURSE_BRIEF}")
    print(divider)
    print("Designing course pack (skill)...\n")

    execution = await agent.async_run_skills_task(
        COURSE_BRIEF,
        mode="required",
        effort="normal",
        output={
            "course_title": (str, "Course title", True),
            "course_overview": (str, "Short course overview", True),
            "teacher_guide_sections": (
                [{"heading": (str, "Section heading", True), "points": ([str], "3-5 teaching points", True)}],
                "Teacher guide sections", True,
            ),
            "student_handout_paragraphs": ([str], "4-6 learner-facing paragraphs", True),
            "lesson_slides": (
                [{"title": (str, "Slide title", True), "bullets": ([str], "3-5 bullets", True)}],
                "6-10 lesson slides", True,
            ),
            "vocabulary": (
                [{"term": (str, "Term", True), "definition": (str, "Definition", True)}],
                "8-12 vocabulary pairs", True,
            ),
            "rubric_criteria": (
                [{"criterion": (str, "Criterion", True), "descriptor": (str, "Performance descriptor", True)}],
                "4-6 rubric criteria", True,
            ),
        },
    )

    print(f"skill status: {execution.status}")
    if execution.status != "success":
        print("output:", execution.output)
        return

    course = execution.output or {}
    out_dir = Path(tempfile.mkdtemp(prefix="agently_course_"))
    written = write_artifacts(out_dir, course)

    print(f"\n  course title: {course.get('course_title', '—')}")
    print(f"  teacher guide sections: {len(course.get('teacher_guide_sections', []) or [])}")
    print(f"  lesson slides: {len(course.get('lesson_slides', []) or [])}")
    print(f"  vocabulary: {len(course.get('vocabulary', []) or [])}")
    print("\n  artifacts written:")
    for name, path in written.items():
        print(f"    · {name}: {path}")

    print(f"\nskill status: {execution.status}")
    print(f"course title: {course.get('course_title', '—')}")
    print(f"artifacts written: {len(written)}")


if __name__ == "__main__":
    asyncio.run(main())
